"""
docs/bag-darga-tanilsuulga.md → docs/bag-darga-tanilsuulga.pptx

Хөдөөгийн багийн даргад зориулсан танилцуулгыг PowerPoint руу хувиргана.
Том үсэг, тод өнгө, олон зурагны байрлал — малчин/ахмадад зориулсан загвар.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

MOCKUP_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..',
                                          'docs', 'mockups'))


def mockup(name):
    """mockups/<name> руу абсолют зам буцаах"""
    return os.path.join(MOCKUP_DIR, name)

# ----- Өнгө (theme) -----
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)  # Гол өнгө — хээр тал
GREEN_MID = RGBColor(0x2E, 0x7D, 0x32)
GOLD = RGBColor(0xFF, 0xB8, 0x00)         # Чухал зүйл онцлох
CREAM = RGBColor(0xFF, 0xF8, 0xE1)        # Цайвар фон
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x55, 0x55, 0x55)
RED_SOFT = RGBColor(0xC6, 0x28, 0x28)

# ----- Хэмжээс (16:9, 13.33 x 7.5 inch) -----
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=CREAM):
    """Слайдын дэвсгэр өнгө"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=24, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             font='Arial'):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header_bar(slide, title, subtitle=None):
    """Дээд талд ногоон тууз — гарчигтай"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_DARK
    bar.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             title, size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
                 subtitle, size=14, color=CREAM)


def add_image_placeholder(slide, left, top, width, height, label,
                          image_name=None):
    """Зурагны байрлал. image_name өгсөн бол жинхэнэ зургийг оруулна.
    Эс бөгөөс алтлаг хүрээтэй текст placeholder."""
    if image_name and os.path.exists(mockup(image_name)):
        return add_phone_screenshot(slide, left, top, width, height,
                                    image_name)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = GOLD
    box.line.width = Pt(3)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run1 = p.add_run()
    run1.text = "📷 ЗУРАГ"
    run1.font.name = 'Arial'
    run1.font.size = Pt(24)
    run1.font.bold = True
    run1.font.color.rgb = GOLD
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "\n" + label
    r2.font.name = 'Arial'
    r2.font.size = Pt(14)
    r2.font.color.rgb = GRAY
    return box


def add_phone_screenshot(slide, area_left, area_top, area_w, area_h,
                         image_name):
    """Утасны frame дотор screenshot оруулах.
    area_left/top/w/h — зургийг байрлуулах нийт талбар.
    Зураг 720x1280 (9:16) тул талбарын дотор төвлөрөх ба
    хар утасны хүрээтэй харагдана."""
    # Утасны зургийн ratio 720/1280
    ratio = 720 / 1280
    # Хамгийн их хэмжээ — area-д багтахаар
    if (area_w / area_h) > ratio:
        # area нь өргөн — өндөрт нь тулгуурлана
        img_h = area_h
        img_w = int(img_h * ratio)
    else:
        img_w = area_w
        img_h = int(img_w / ratio)
    img_left = area_left + (area_w - img_w) // 2
    img_top = area_top + (area_h - img_h) // 2
    # Хар утасны frame (зургаас 6% том)
    pad = int(img_w * 0.04)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   img_left - pad, img_top - pad,
                                   img_w + pad * 2, img_h + pad * 2)
    frame.fill.solid()
    frame.fill.fore_color.rgb = BLACK
    frame.line.color.rgb = GRAY
    frame.line.width = Pt(2)
    frame.adjustments[0] = 0.05
    # Жинхэнэ зураг
    pic = slide.shapes.add_picture(mockup(image_name),
                                   img_left, img_top, img_w, img_h)
    return pic


def add_bullet_box(slide, left, top, width, height, items, *, size=18,
                   color=BLACK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "• " + item
        run.font.name = 'Arial'
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def add_advantage_card(slide, left, top, width, height, title, body):
    """Алтан хүрээтэй давуу талын карт"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = GOLD
    card.line.width = Pt(2)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "✅ " + title
    r.font.name = 'Arial'
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = GREEN_DARK
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = 'Arial'
    r2.font.size = Pt(13)
    r2.font.color.rgb = BLACK
    return card


# ============================================================
# СЛАЙД 1 — НҮҮР
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, GREEN_DARK)

# Том лого/гарчиг
add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.2),
         "🐑 МАЛЧИН АПП", size=72, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.8),
         "Багийн даргад зориулсан танилцуулга",
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
         "«Малчны өглөөний дэвтэр, өдрийн зөвлөх, оройн тооцоо»",
         size=20, color=CREAM, align=PP_ALIGN.CENTER)

# Доод цагираг
ribbon = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.7), SLIDE_W, Inches(0.8))
ribbon.fill.solid()
ribbon.fill.fore_color.rgb = GOLD
ribbon.line.fill.background()
add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5),
         "Хөдөө сум, баг, малчид — нэг утсан дээр",
         size=18, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ============================================================
# СЛАЙД 2 — ЭНЭ АПП ГЭЖ ЮУ ВЭ?
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "1. ЭНЭ АПП ГЭЖ ЮУ ВЭ?",
               "Богино үгээр: «Өглөө бүр нээдэг малчны дэвтэр»")

add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.8),
         "Малчин хүн утсаа авмагц өнөөдрийн бүх мэдээллээ нэг дороос харна:",
         size=20, color=BLACK)

items = [
    "Малын тоо толгой, бэлчээр, өвөлжөө хаваржаа",
    "Цаг агаар, зуд, цасны мэдээ — өөрийн нутгаар",
    "Сум, багийн зарлал — алгасахгүй",
    "Алдсан мал, олдсон малыг 2 минутад зарлах",
    "Өвс тэжээл, унаа тээвэр, ажиллах хүчний зар",
    "Малын эмчтэй уулзах цаг захиалах",
    "Ахмад малчдын ухаан, зөвлөгөө",
    "Орлого зарлагын богино тэмдэглэл",
]
add_bullet_box(s, Inches(0.8), Inches(2.2), Inches(12), Inches(5),
               items, size=20)


# ============================================================
# СЛАЙД 3 — ЯАГААД МАЛЧДАД ХЭРЭГТЭЙ ВЭ
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "2. ЯАГААД МАЛЧДАД ХЭРЭГТЭЙ ВЭ?",
               "Өнөөдөр малчин үүнийг бүгдийг тусдаа газраас хайдаг")

# Хүснэгт — 3 багана
headers = ["Юу хайдаг вэ", "Хаанаас хайдаг вэ", "Асуудал"]
rows = [
    ["Цаг агаар", "Зурагт, радио", "Өөрийн нутгийн биш"],
    ["Сумын мэдэгдэл", "Утсан мессеж", "Дараа алддаг"],
    ["Алдсан мал", "Айлд эргэх", "Удаан, хол хүртдэггүй"],
    ["Өвс тэжээлийн үнэ", "Сумын төв", "Хол явахаас өмнө мэдэхгүй"],
    ["Зөвлөгөө", "Ахмад айл", "Ахмад нь холдсон"],
    ["Малын тоо", "Толгойдоо", "Хэдэн жилд мартагддаг"],
]

# Хүснэгт байрлуулах
left, top = Inches(0.5), Inches(1.4)
col_w = [Inches(3.7), Inches(4.2), Inches(4.7)]
row_h = Inches(0.55)

# Толгой
x = left
for i, h in enumerate(headers):
    cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, col_w[i], row_h)
    cell.fill.solid()
    cell.fill.fore_color.rgb = GREEN_DARK
    cell.line.color.rgb = WHITE
    tf = cell.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = h
    r.font.name = 'Arial'
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE
    x += col_w[i]

# Мөрнүүд
y = top + row_h
for ri, row in enumerate(rows):
    x = left
    bg_color = WHITE if ri % 2 == 0 else CREAM
    for i, val in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w[i], row_h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        cell.line.color.rgb = GRAY
        tf = cell.text_frame
        tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = val
        r.font.name = 'Arial'
        r.font.size = Pt(14)
        r.font.color.rgb = BLACK
        x += col_w[i]
    y += row_h

# Дүгнэлт мөр
add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.9),
         "👉 Энэ апп бүгдийг нэг газар цуглуулсан. Сүлжээгүй ч ажиллана.",
         size=20, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ============================================================
# СЛАЙДЫН ҮНДСЭН ТУСЛАГЧ — Хуудас бүрд нэг слайд
# ============================================================
def add_screen_slide(num, title, image_label, description, advantages,
                     image_name=None):
    """Аппын нэг хуудсыг танилцуулах слайд:
    зүүн талд зурагны байрлал, баруун талд тайлбар + давуу тал"""
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header_bar(s, f"3.{num}  {title}", "Аппын хуудас")

    # Зүүн тал — зурагны байрлал
    add_image_placeholder(s, Inches(0.5), Inches(1.4),
                          Inches(5.2), Inches(5.7), image_label,
                          image_name=image_name)

    # Баруун тал — тайлбар
    add_text(s, Inches(6.0), Inches(1.4), Inches(7), Inches(0.5),
             "📝 Тайлбар", size=20, bold=True, color=GREEN_DARK)
    add_bullet_box(s, Inches(6.0), Inches(1.9), Inches(7), Inches(3.0),
                   description, size=15)

    # Баруун доод — давуу тал
    add_text(s, Inches(6.0), Inches(5.0), Inches(7), Inches(0.5),
             "✅ Давуу тал", size=20, bold=True, color=GOLD)
    add_bullet_box(s, Inches(6.0), Inches(5.5), Inches(7), Inches(2.0),
                   advantages, size=14)

    return s


# ============================================================
# СЛАЙД 4 — Нэвтрэх (3.1)
# ============================================================
add_screen_slide(
    "1. НЭВТРЭХ",
    "Утасны дугаараар нэвтрэх",
    "Утасны дугаар оруулах хуудас\n\nДараа нь мессежээр нууц код ирнэ",
    [
        "Зөвхөн утасны дугаараа бичнэ",
        "Мессежээр нууц код ирнэ",
        "Кодоо бичээд орно",
        "Имэйл, нууц үг, нэмэлт зүйл байхгүй",
    ],
    [
        "Хөгшин малчин ч ойлгож, өөрөө бүртгүүлж чадна",
        "Имэйл бий болгох шаардлагагүй",
    ],
    image_name="01-phone.png",
)


# ============================================================
# СЛАЙД 5 — 10 алхамт бүртгэл (3.2)
# ============================================================
add_screen_slide(
    "2. ӨӨРИЙГӨӨ ТАНИУЛАХ",
    "10 алхамт нэг удаагийн бүртгэл",
    "Овог нэр, аймаг сум баг сонгох\n\n4 улирлын байршлыг тэмдэглэх\n\nМалын тоог оруулах",
    [
        "Овог, нэр",
        "Аймаг, сум, баг",
        "Та хэн бэ? (малчин/багийн дарга/...)",
        "4 улирлын байршил — өвөлжөө, хаваржаа, зуслан, намаржаа",
        "Одоо хаана байгаа",
        "Оторлох боломжтой газар",
        "Малын тоо толгой (5 төрөл)",
    ],
    [
        "Нэг л удаа бөглөнө — дараа нь дахин асуухгүй",
        "Апп өөрөө таны нутаг, малыг таниад мэдээлэл өгнө",
    ],
    image_name="02-role.png",
)


# ============================================================
# СЛАЙД 6 — Нүүр хуудас (3.3) — ТУСГАЙ ХЭВ ЗАГВАР
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "3.3  НҮҮР ХУУДАС — «ӨНӨӨДӨР»",
               "Өглөө апп нээмэгц гарах 9 карт")

# Зүүн дээр — зурагны байрлал
add_image_placeholder(s, Inches(0.5), Inches(1.4),
                      Inches(4.5), Inches(5.7),
                      "Нүүр хуудас\n\nЦаг агаар, өнөөдөр хийх ажил, мэдэгдэл",
                      image_name="03-home.png")

# Баруун — 9 картын жагсаалт
cards = [
    ("🌤", "Өнөөдрийн цаг агаар", "Өвөлжөөны байршлаар"),
    ("⚠️", "Өнөөдрийн эрсдэл", "Салхи, хүйтэн, цас, зуд"),
    ("✅", "Өнөөдөр хийх 3 ажил", "Ус, түлш, эрүүл мэнд"),
    ("🚚", "Нүүх / оторлох зөвлөгөө", "Өвс хүрэхгүй бол хаашаа"),
    ("🩺", "Малын эрүүл мэндийн дохио", "Шүлхий, таршаа г.м."),
    ("📢", "Сумын шинэ мэдэгдэл", "Багийн дарга юу мэдэгдсэн"),
    ("🧓", "Ахмадын 1 зөвлөгөө", "Өдөрт нэг сургаалт үг"),
    ("💰", "Зах зээлийн товч үнэ", "Өвс, мах, ноос"),
    ("📦", "Ойролцоох хэрэгтэй зар", "Унаа, ажиллах хүн"),
]
left2 = Inches(5.3)
top2 = Inches(1.4)
card_w = Inches(2.55)
card_h = Inches(1.8)
gap = Inches(0.1)
for idx, (emoji, title, sub) in enumerate(cards):
    col = idx % 3
    row = idx // 3
    cl = left2 + (card_w + gap) * col
    ct = top2 + (card_h + gap) * row
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, ct, card_w, card_h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = GREEN_MID
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = emoji
    r.font.size = Pt(28)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = title
    r2.font.name = 'Arial'
    r2.font.size = Pt(11)
    r2.font.bold = True
    r2.font.color.rgb = GREEN_DARK
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = sub
    r3.font.name = 'Arial'
    r3.font.size = Pt(9)
    r3.font.color.rgb = GRAY


# ============================================================
# СЛАЙД 7 — Мал бүртгэл
# ============================================================
add_screen_slide(
    "4. МАЛ БҮРТГЭЛ",
    "5 төрлийн мал — адуу, үхэр, хонь, ямаа, тэмээ",
    "Малын тоо толгой\n\nНэмэх, хасах товч энгийн",
    [
        "5 төрлийн малыг тус тусад нь бүртгэнэ",
        "Нэмэх, хасах товч энгийн",
        "Төл гарах, хорогдох, худалдсан, авсан — бичнэ",
        "Жилийн эцсийн тооллогод бэлэн болно",
    ],
    [
        "«Хэдэн толгой малтай вэ?» гэхэд 3 секундэд хариулна",
        "Толгойдоо барих, хуучин дэвтэрт хайх шаардлагагүй",
    ],
    image_name="04-livestock.png",
)


# ============================================================
# СЛАЙД 8 — Бэлчээр / Газрын зураг
# ============================================================
add_screen_slide(
    "5. БЭЛЧЭЭР, ГАЗРЫН ЗУРАГ",
    "4 улирлын бууц + одоогийн байршил",
    "Газрын зураг\n\nӨвөлжөө, хаваржаа, зуслан, намаржаа",
    [
        "4 улирлын өөрийн бууцыг тэмдэглэнэ",
        "Одоо мал хаана байгаа газрын зурагт цэгээр",
        "Оторлох боломжтой газар нэмэх",
    ],
    [
        "Яг газрыг бусдад харуулахгүй — зөвхөн «энэ багт» түвшинд",
        "Хувийн нууц хадгалагдана",
    ],
    image_name="05-pasture.png",
)


# ============================================================
# СЛАЙД 9 — Зөвлөгөө (Ухаалаг туслах)
# ============================================================
add_screen_slide(
    "6. ЗӨВЛӨГӨӨ — УХААЛАГ ТУСЛАХ",
    "Малчны асуултанд богино, тодорхой хариулт",
    "Зөвлөгөө асуух хуудас\n\nХариултын карт",
    [
        "«Хонь толгой эргэвэл яах вэ?»",
        "«Зуд болох гэж байна, хаашаа нүүх вэ?»",
        "Хариу: одоо юу хийх → яагаад → 3-7 алхам → анхаарах → эрсдэл → хэнд мэдэгдэх",
    ],
    [
        "Урт өгүүлэл биш — богино тодорхой алхмууд",
        "Малчинд шууд хэрэгцээтэй формат",
    ],
    image_name="06-advisory.png",
)


# ============================================================
# СЛАЙД 10 — Ахмадын ухаан
# ============================================================
add_screen_slide(
    "7. АХМАДЫН УХААН",
    "Ахмад малчдын дуу хоолой, видео, зураг",
    "Хөгшин малчны видео карт\n\nАудио сонсох товч",
    [
        "Ахмадын дуу хоолой, видео, зураг",
        "«Зуны бэлчээр», «Тулга бариулах», «Чонын дайралтаас»",
        "Хөгшин хүмүүс өөрсдөө оруулна",
    ],
    [
        "Хуучин ёс заншил, малын ухаан залуу үед дамжина",
        "Ахмадууд өөрсдөө хүндлэгдэж байгаагаа мэдэрнэ",
    ],
    image_name="07-elder.png",
)


# ============================================================
# СЛАЙД 11 — Алдсан/Олдсон мал
# ============================================================
add_screen_slide(
    "8. АЛДСАН / ОЛДСОН МАЛ",
    "2 минутад зарлах тусгай хуудас",
    "Алдсан мал зарлах хуудас\n\nОйролцоох олдсон малын жагсаалт",
    [
        "АЛДСАН: зураг авах → им тамга бичих → багт мэдэгдэх",
        "ОЛДСОН: зураг авах → нутгийн зар → эзэнтэй холбогдох",
    ],
    [
        "Урьд сум сум руу яриулж байсан зүйл одоо 2 минутад",
        "Хол хүртэлх айлуудад нэг дор хүрнэ",
    ],
    image_name="08-lost-found.png",
)


# ============================================================
# СЛАЙД 12 — Зах зээл
# ============================================================
add_screen_slide(
    "9. ЗАХ ЗЭЭЛ — ЗАР ЗАХИАЛГА",
    "Малчны амьдралтай холбоотой зар",
    "Зар жагсаалт\n\nӨвс тэжээл, унаа, мал, ажиллах хүн",
    [
        "Өвс тэжээл, хужир шорлог",
        "Тээвэр, унаа",
        "Мал зарах, авах",
        "Ажиллах хүн, түрээс",
        "Хоршооны бараа",
    ],
    [
        "Хол явалгүй сум, аймгийн мэдээлэл харна",
        "Үнэ ил тод, дундын зуучлагч хэрэггүй",
    ],
    image_name="09-market.png",
)


# ============================================================
# СЛАЙД 13 — Мал эмч
# ============================================================
add_screen_slide(
    "10. МАЛ ЭМЧТЭЙ ХОЛБОГДОХ",
    "Цаг захиалах, асуулт илгээх",
    "Мал эмчийн жагсаалт\n\nЦаг захиалах хуудас",
    [
        "Сумын мал эмчийн утас, хаяг",
        "Цаг захиалах",
        "Малын зураг илгээх, асуулт бичих",
        "Эмийн зөвлөгөө",
    ],
    [
        "«Эмч сум руу очно уу?» гэж хүлээх биш",
        "Апп дотроос шууд ярина",
    ],
    image_name="10-vet.png",
)


# ============================================================
# СЛАЙД 14 — Сумын мэдэгдэл
# ============================================================
add_screen_slide(
    "11. СУМ, БАГИЙН МЭДЭГДЭЛ",
    "Зарлал, тарилгын хуваарь, хурлын мэдэгдэл",
    "Сумын зарлал жагсаалт",
    [
        "Сумын засаг даргын зарлал",
        "Багийн даргын мэдэгдэл",
        "Тооллого, тарилгын хуваарь",
        "Хурал, цуглааны мэдэгдэл",
    ],
    [
        "«Сонссонгүй, мэдсэнгүй» гэх асуудал тэглэдэг",
        "Хэн уншсан, хэн уншаагүй харагдана",
    ],
    image_name="11-news.png",
)


# ============================================================
# СЛАЙД 15 — Өрх, санхүү
# ============================================================
add_screen_slide(
    "12. ӨРХ, САНХҮҮ",
    "Орлого зарлага, өрхийн гишүүд",
    "Орлого зарлагын тэмдэглэл\n\nӨрхийн гишүүд",
    [
        "Өрхийн гишүүд (эхнэр, хүүхэд, хадам)",
        "Орлого, зарлагын богино тэмдэглэл",
        "ЭМД, НДШ-ийн мэдээлэл",
        "Хүүхдийн сургуулийн дэмжлэг",
    ],
    [
        "Орлого зарлагаа толгойдоо барилгүй харна",
        "Жилийн тооцоо хийхэд хялбар",
    ],
    image_name="12-household.png",
)


# ============================================================
# СЛАЙД 16 — Багц, төлбөр
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "3.13  ПРОФАЙЛ, БАГЦ, ТӨЛБӨР", "5 багц — үнэгүйгээр эхэлж болно")

packages = [
    ("ҮНЭГҮЙ", "Хэн ч", "Үндсэн боломжууд", GREEN_MID),
    ("ПРЕМИУМ МАЛЧИН", "Малчин өөрөө", "Илүү гүн зөвлөгөө,\nэрсдэлийн анхааруулга", GOLD),
    ("ХОРШООНЫ БАГЦ", "Хоршоо", "Олон гишүүн нэгтгэх", GREEN_DARK),
    ("СУМЫН БАГЦ", "Сумын засаг", "Бүх малчдыг харах\nсамбар", GREEN_DARK),
    ("ҮЙЛЧИЛГЭЭ", "Мал эмч, тээвэр", "Зар, цаг захиалга", GOLD),
]

card_w2 = Inches(2.4)
card_h2 = Inches(3.5)
left3 = Inches(0.6)
top3 = Inches(1.5)
gap2 = Inches(0.15)

for i, (name, who, what, color) in enumerate(packages):
    cl = left3 + (card_w2 + gap2) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, top3, card_w2, card_h2)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(3)
    # Дээд тал — гарчиг
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cl, top3, card_w2, Inches(0.7))
    head.fill.solid()
    head.fill.fore_color.rgb = color
    head.line.fill.background()
    tf = head.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = name
    r.font.name = 'Arial'
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE
    # Биеийн текст
    body_box = s.shapes.add_textbox(cl, top3 + Inches(0.85),
                                    card_w2, Inches(2.5))
    btf = body_box.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = "Хэн авах:"
    br.font.name = 'Arial'
    br.font.size = Pt(11)
    br.font.color.rgb = GRAY
    bp2 = btf.add_paragraph()
    bp2.alignment = PP_ALIGN.CENTER
    br2 = bp2.add_run()
    br2.text = who
    br2.font.name = 'Arial'
    br2.font.size = Pt(13)
    br2.font.bold = True
    br2.font.color.rgb = BLACK
    bp3 = btf.add_paragraph()
    bp3.alignment = PP_ALIGN.CENTER
    bp3.space_before = Pt(8)
    br3 = bp3.add_run()
    br3.text = "\nЮу багтсан:"
    br3.font.name = 'Arial'
    br3.font.size = Pt(11)
    br3.font.color.rgb = GRAY
    bp4 = btf.add_paragraph()
    bp4.alignment = PP_ALIGN.CENTER
    br4 = bp4.add_run()
    br4.text = what
    br4.font.name = 'Arial'
    br4.font.size = Pt(12)
    br4.font.color.rgb = BLACK

add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.6),
         "💡 Үнэгүй хувилбараар л өдөр тутмын ажил гүйцэт бүтнэ",
         size=20, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ============================================================
# СЛАЙД 17 — БАГИЙН ДАРГЫН ТУСГАЙ ХУУДАС (4)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "4. БАГИЙН ДАРГЫН ТУСГАЙ САМБАР",
               "Зөвхөн багийн даргад нээгдэх боломжууд")

add_image_placeholder(s, Inches(0.5), Inches(1.4),
                      Inches(5), Inches(5.7),
                      "Багийн даргын самбар\n\nӨрхүүдийн жагсаалт, эрсдэл,\nмэдэгдэл явуулах товч",
                      image_name="14-bag-dashboard.png")

# 4 ач холбогдол
benefits = [
    ("🏠", "Бүх өрхийг нэг дэлгэцэн дээр",
     "Хэдэн өрх, хэчнээн мал. Хэн нь өвөлжөө руу нүүсэн, хэн нь оторт явсан"),
    ("📢", "Нэг товчоор бүгдэд мэдэгдэл",
     "Гарчиг, агуулга бичээд илгээх. Хэн уншсан, хэн уншаагүй харагдана"),
    ("⚠️", "Эрсдэлтэй өрхүүд тодрох",
     "Цастай нутаг, өвс хүрэхгүй, хол оторт яваа. Шууд утсаар холбогдох"),
    ("📊", "Тайлан өөрөө гарна",
     "Сар, улирлын тайлан. Сумын засаг даргад илгээх товчтой"),
]

bx = Inches(5.8)
by = Inches(1.5)
bw = Inches(7.2)
bh = Inches(1.3)
gap3 = Inches(0.15)

for i, (emo, title, body) in enumerate(benefits):
    bt = by + (bh + gap3) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, bt, bw, bh)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = GREEN_DARK
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{emo}  {title}"
    r.font.name = 'Arial'
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = GREEN_DARK
    p2 = tf.add_paragraph()
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = 'Arial'
    r2.font.size = Pt(12)
    r2.font.color.rgb = BLACK


# ============================================================
# СЛАЙД 18 — ДАВУУ ТАЛУУД (хураангуй)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "5. ДАВУУ ТАЛУУД (ХУРААНГУЙ)",
               "Малчин, багийн дарга, сумын удирдлага — гурвуулангаа хожино")

# 3 багана
columns = [
    ("👨‍🌾  МАЛЧДАД", GREEN_DARK, [
        "Бүх мэдээллээ нэг газраас",
        "Сүлжээгүй ч ажиллана",
        "Хол явалгүй өвсний үнэ",
        "Хөгшин ч ойлгох товч",
        "Алдсан мал 2 минутад",
        "Малын тоо толгойдоо хадгалгүй",
    ]),
    ("👨‍💼  БАГИЙН ДАРГАД", GOLD, [
        "Бүх өрхөө нэг самбараас",
        "Нэг товчоор мэдэгдэл",
        "Эрсдэлтэй өрхөө хурдан",
        "Тайлан цаг алдалгүй",
    ]),
    ("🏛  СУМЫН УДИРДЛАГАД", GREEN_MID, [
        "Бүх багийн мэдээлэл",
        "Өрхүүдтэй шууд",
        "Тайлан, статистик гарт",
    ]),
]

cw = Inches(4.1)
ch = Inches(5.5)
cleft = Inches(0.5)
ctop = Inches(1.5)
cgap = Inches(0.2)

for i, (title, color, items) in enumerate(columns):
    cl = cleft + (cw + cgap) * i
    # Толгой
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cl, ctop, cw, Inches(0.7))
    head.fill.solid()
    head.fill.fore_color.rgb = color
    head.line.fill.background()
    tf = head.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = 'Arial'
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    # Бие
    body = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cl, ctop + Inches(0.7),
                              cw, ch - Inches(0.7))
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    body.line.color.rgb = color
    body.line.width = Pt(2)
    btf = body.text_frame
    btf.word_wrap = True
    btf.margin_left = Inches(0.25)
    btf.margin_top = Inches(0.2)
    for j, item in enumerate(items):
        if j == 0:
            bp = btf.paragraphs[0]
        else:
            bp = btf.add_paragraph()
        bp.space_after = Pt(8)
        br = bp.add_run()
        br.text = "✓ " + item
        br.font.name = 'Arial'
        br.font.size = Pt(13)
        br.font.color.rgb = BLACK


# ============================================================
# СЛАЙД 19 — ЯАГААД ЗААВАЛ (7 шалтгаан)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "6. ЯАГААД ЗААВАЛ ЭНЭ АППЫГ АШИГЛАХ ВЭ?",
               "7 чухал шалтгаан")

reasons = [
    ("1", "Цаг хэмнэнэ", "Утсаар яриа, дэвтэр бичих, сум руу явах хугацаа богиносно"),
    ("2", "Алдагдал багасна", "Цаг агаар, өвчний анхааруулгыг урьдчилан мэднэ"),
    ("3", "Үнэ ил тод", "Дундын зуучлагч «бага үнэлэх» эрсдэл багасна"),
    ("4", "Хүн бүр ижил мэдээлэл", "Багийн дарга, малчин, сумын дарга нэг тоог хардаг"),
    ("5", "Хөгшин залуу хоёр ойлгоно", "Хэцүү технологи биш"),
    ("6", "Хувийн нууц хадгалагдана", "Яг координат бусдад харагдахгүй"),
    ("7", "Үнэгүй эхэлж болно", "0 төгрөгөөр өдөр тутмын ажил бүтнэ"),
]

# 2 багана 4 мөр layout
left5 = Inches(0.5)
top5 = Inches(1.4)
rw = Inches(6.2)
rh = Inches(1.4)
rgap_x = Inches(0.2)
rgap_y = Inches(0.1)

for i, (num, title, body) in enumerate(reasons):
    col = i % 2
    row = i // 2
    rl = left5 + (rw + rgap_x) * col
    rt = top5 + (rh + rgap_y) * row
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rl, rt, rw, rh)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = GOLD
    box.line.width = Pt(2)
    # Тоо
    num_box = s.shapes.add_shape(MSO_SHAPE.OVAL, rl + Inches(0.15),
                                 rt + Inches(0.3), Inches(0.8), Inches(0.8))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = GOLD
    num_box.line.fill.background()
    ntf = num_box.text_frame
    ntf.margin_left = 0
    ntf.margin_right = 0
    ntf.margin_top = 0
    ntf.margin_bottom = 0
    np = ntf.paragraphs[0]
    np.alignment = PP_ALIGN.CENTER
    nr = np.add_run()
    nr.text = num
    nr.font.name = 'Arial'
    nr.font.size = Pt(28)
    nr.font.bold = True
    nr.font.color.rgb = WHITE
    # Текст
    add_text(s, rl + Inches(1.1), rt + Inches(0.15),
             rw - Inches(1.2), Inches(0.55), title,
             size=18, bold=True, color=GREEN_DARK)
    add_text(s, rl + Inches(1.1), rt + Inches(0.7),
             rw - Inches(1.2), Inches(0.6), body,
             size=12, color=BLACK)


# ============================================================
# СЛАЙД 20 — НЭВТРҮҮЛЭХ 5 АЛХАМ
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "7. БАГТАА ЯАЖ НЭВТРҮҮЛЭХ ВЭ?",
               "Багийн дарга та өөрөө хэрхэн эхлүүлэх")

steps = [
    ("1", "БАГИЙН ХУРАЛ ДЭЭР",
     "Нэг утсаар үзүүлж тайлбарлах. Бүгд харж, асуухаар"),
    ("2", "ИДЭВХТЭЙ АЙЛУУДААР ЭХЭЛ",
     "3-5 идэвхтэй малчин эхлээд туршиж үзнэ"),
    ("3", "АМЖИЛТЫГ ҮЗҮҮЛ",
     "Жишээ нь алдсан мал хурдан олдсон туршлагыг бусдад"),
    ("4", "ХУУЧИН АРГЫГ ҮЛДЭЭ",
     "Хуучин арга барилаа хэвээр үлдээгээд дээр нь нэмэлт хэрэгсэл"),
    ("5", "1-2 САРЫН ДОТОР",
     "Бүх айлд тарж ашиглагдана"),
]

slt = Inches(1.0)
slx_start = Inches(0.4)
sw = Inches(2.5)
sh = Inches(5)
sgap = Inches(0.15)

for i, (num, title, body) in enumerate(steps):
    sl = slx_start + (sw + sgap) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sl, Inches(1.5), sw, sh)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = GREEN_DARK
    box.line.width = Pt(2)
    # Дугаар тойрог
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                sl + sw/2 - Inches(0.6),
                                Inches(1.7), Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN_DARK
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.margin_left = 0
    ctf.margin_right = 0
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = num
    cr.font.name = 'Arial'
    cr.font.size = Pt(40)
    cr.font.bold = True
    cr.font.color.rgb = WHITE
    # Гарчиг
    add_text(s, sl + Inches(0.1), Inches(3.1),
             sw - Inches(0.2), Inches(1.0),
             title, size=14, bold=True, color=GREEN_DARK,
             align=PP_ALIGN.CENTER)
    # Тайлбар
    add_text(s, sl + Inches(0.15), Inches(4.2),
             sw - Inches(0.3), Inches(2.0),
             body, size=12, color=BLACK, align=PP_ALIGN.CENTER)


# ============================================================
# СЛАЙД 21 — ХУРААНГУЙ / БАЯРЛАЛАА
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, GREEN_DARK)

add_text(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.8),
         "🐑 МАЛЧИН АПП", size=44, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER)

# Том дөрвөлжин дотор хураангуй
quote_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1), Inches(2),
                               Inches(11.3), Inches(2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = CREAM
quote_box.line.color.rgb = GOLD
quote_box.line.width = Pt(3)
qtf = quote_box.text_frame
qtf.word_wrap = True
qtf.margin_left = Inches(0.3)
qtf.margin_top = Inches(0.3)
qp = qtf.paragraphs[0]
qp.alignment = PP_ALIGN.CENTER
qr = qp.add_run()
qr.text = "«Малчны өглөөний дэвтэр,"
qr.font.name = 'Arial'
qr.font.size = Pt(28)
qr.font.bold = True
qr.font.color.rgb = GREEN_DARK
qp2 = qtf.add_paragraph()
qp2.alignment = PP_ALIGN.CENTER
qr2 = qp2.add_run()
qr2.text = "өдрийн зөвлөх, оройн тооцоо"
qr2.font.name = 'Arial'
qr2.font.size = Pt(28)
qr2.font.bold = True
qr2.font.color.rgb = GREEN_DARK
qp3 = qtf.add_paragraph()
qp3.alignment = PP_ALIGN.CENTER
qr3 = qp3.add_run()
qr3.text = "— нэг утсан дээр»"
qr3.font.name = 'Arial'
qr3.font.size = Pt(28)
qr3.font.bold = True
qr3.font.color.rgb = GREEN_DARK

# Багийн даргын ач тус
benefits_final = [
    "🏠  Багийн БҮХ ӨРХТЭЙ холбогдсон хэвээр",
    "📊  Тайлан илгээх ажил ХӨНГӨВЧЛӨНӨ",
    "⚠️  Эрсдэл ирэхэд ЦАГ АЛДАЛГҮЙ сэрэмжлүүлнэ",
    "🏛  Сумын удирдлагатай БЭЛЭН МЭДЭЭЛЛЭЭР",
]

for i, item in enumerate(benefits_final):
    add_text(s, Inches(2), Inches(4.4) + Inches(0.5) * i,
             Inches(9.5), Inches(0.5),
             item, size=18, bold=True, color=WHITE)

add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
         "БАЯРЛАЛАА  •  АСУУЛТ БАЙВАЛ АПП ДОТРООС «ТУСЛАМЖ» ХЭСГИЙГ ҮЗНЭ ҮҮ",
         size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ============================================================
# Хадгалах
# ============================================================
out = os.path.join(os.path.dirname(__file__), '..', 'docs',
                   'bag-darga-tanilsuulga.pptx')
out = os.path.abspath(out)
prs.save(out)
print(f"OK saved: {out}")
print(f"   Slides: {len(prs.slides)}")
